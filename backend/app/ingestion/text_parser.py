"""
Universal text parser for plain-text formats and PowerPoint.
Handles: .txt .md .csv .json .log .py .pptx
"""
import os
from typing import List, Dict

PARAGRAPHS_PER_PAGE = 12


def parse_text_file(file_path: str) -> List[Dict]:
    """Parse a plain-text file into page-like sections."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    lines = [l for l in content.split("\n") if l.strip()]
    if not lines:
        return [{"page_number": 1, "text": ""}]

    pages = []
    for i in range(0, len(lines), PARAGRAPHS_PER_PAGE):
        chunk = lines[i : i + PARAGRAPHS_PER_PAGE]
        pages.append({
            "page_number": (i // PARAGRAPHS_PER_PAGE) + 1,
            "text": "\n".join(chunk),
            "section_kind": "section",
        })
    return pages


def parse_pptx(file_path: str) -> List[Dict]:
    """Parse a PowerPoint file — one 'page' per slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    pages = []

    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        slide_text = "\n".join(texts)
        if slide_text.strip():
            pages.append({"page_number": idx, "text": slide_text})

    if not pages:
        return [{"page_number": 1, "text": ""}]
    return pages
